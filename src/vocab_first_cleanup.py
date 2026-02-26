#!/usr/bin/env python3
"""
Vocab-First Cleanup Script

Brings the existing database to a vocab-first state in four steps:

  Step 1 — Delete noise
      Remove all occurrences with validation_status IN
      ('potential_noise', 'high_priority_review').
      Clean up orphan concepts afterwards.

  Step 2 — Promote confirmed_with_flag
      Set validation_status = 'confirmed' for all confirmed_with_flag rows.
      Vocab list membership overrides Stage 1 noise flags.

  Step 3 — Update chapters from vocab lists
      For surviving occurrences with NULL/empty chapter, populate from vocab
      list chapter structure. Where DB chapter number conflicts with vocab
      list chapter number, log the conflict but do not overwrite full chapter
      titles with bare numbers.

  Step 4 — Recover missed vocab terms
      For each unit, find vocab list terms absent from the DB.
      Search PPTX text (all runs, not just bold) for each missed term.
      Insert a confirmed occurrence (is_introduction=0) for those found.
      Terms not present in the PPTX text at all are skipped — they are
      supplementary vocab not used in the booklet body.

Usage:
    python src/vocab_first_cleanup.py [--dry-run] [--skip-promote]

Arguments:
    --dry-run       Print counts only; do not modify the database.
    --skip-promote  Skip Step 2 (keep confirmed_with_flag for manual review).

Created: 2026-02-26
"""

import argparse
import re
import sqlite3
import sys
from pathlib import Path

# Allow imports from src/ when invoked from project root
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation

# _normalise is a private helper in vocab_validator; used here for fuzzy
# chapter lookup in Step 3. Same project, so acceptable.
from vocab_validator import find_vocab_list, match_term, parse_vocab_docx
from vocab_validator import _normalise  # noqa: PLC2701


# =============================================================================
# DATABASE HELPERS
# =============================================================================

def get_all_units(cursor: sqlite3.Cursor) -> list[dict]:
    """
    Return distinct units from occurrences with one source_path each.

    Created: 2026-02-26
    """
    cursor.execute("""
        SELECT DISTINCT subject, year, term, unit, source_path
        FROM occurrences
        ORDER BY subject, year, term, unit
    """)
    return [
        {
            'subject': r[0], 'year': r[1], 'term': r[2],
            'unit': r[3], 'source_path': r[4]
        }
        for r in cursor.fetchall()
    ]


def get_unit_occurrences(
    cursor: sqlite3.Cursor,
    subject: str,
    year: int,
    term: str,
    unit: str
) -> list[dict]:
    """
    Return all occurrences + concept terms for a unit.

    Created: 2026-02-26
    """
    cursor.execute("""
        SELECT o.occurrence_id, c.term AS concept_term, o.chapter,
               o.slide_number, o.validation_status
        FROM occurrences o
        JOIN concepts c ON o.concept_id = c.concept_id
        WHERE o.subject=? AND o.year=? AND o.term=? AND o.unit=?
    """, (subject, year, term, unit))
    return [
        {
            'occurrence_id': r[0], 'concept_term': r[1], 'chapter': r[2],
            'slide_number': r[3], 'validation_status': r[4]
        }
        for r in cursor.fetchall()
    ]


def cleanup_orphan_concepts(cursor: sqlite3.Cursor) -> int:
    """
    Remove concepts with no remaining occurrences. Returns count deleted.

    Created: 2026-02-26
    """
    cursor.execute("""
        DELETE FROM concepts
        WHERE concept_id NOT IN (SELECT DISTINCT concept_id FROM occurrences)
    """)
    return cursor.rowcount


def get_or_create_concept(cursor: sqlite3.Cursor, term: str) -> int:
    """
    Return concept_id for term; insert if not found.

    Created: 2026-02-26
    """
    cursor.execute("SELECT concept_id FROM concepts WHERE term = ?", (term,))
    row = cursor.fetchone()
    if row:
        return row[0]
    cursor.execute("INSERT INTO concepts (term) VALUES (?)", (term,))
    return cursor.lastrowid


# =============================================================================
# PPTX TEXT SEARCH
# =============================================================================

def search_term_in_pptx(pptx_path: str, term: str) -> dict:
    """
    Search all text (bold and unbolded) in a PPTX for a term.

    Uses word-boundary regex, case-insensitive. Assembles paragraph text
    across all runs to catch terms split at run boundaries.

    Returns dict:
        found         : bool
        slides        : sorted list of 1-based slide numbers
        first_context : paragraph text of first match
        error         : str (only present on failure)

    Created: 2026-02-26
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {'found': False, 'slides': [], 'first_context': '', 'error': str(e)}

    pattern = re.compile(
        r'(?<![a-zA-Z])' + re.escape(term) + r'(?![a-zA-Z])',
        re.IGNORECASE
    )

    matching_slides = []
    first_context = ''

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_matched = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                para_text = ''.join(run.text for run in para.runs)
                if para_text and pattern.search(para_text):
                    slide_matched = True
                    if not first_context:
                        first_context = para_text.strip()
        if slide_matched:
            matching_slides.append(slide_idx)

    return {
        'found': bool(matching_slides),
        'slides': sorted(matching_slides),
        'first_context': first_context
    }


# =============================================================================
# CHAPTER UTILITIES
# =============================================================================

def chapter_number_from_db_string(chapter_str: str | None) -> str | None:
    """
    Extract the leading chapter number from a DB chapter string.

    Examples:
        '1. The Roman Empire' → '1'
        'Chapter 2'           → '2'
        '3'                   → '3'
        None / ''             → None

    Created: 2026-02-26
    """
    if not chapter_str:
        return None
    m = re.match(r'^(?:Chapter\s+)?(\d+)', chapter_str.strip(), re.IGNORECASE)
    return m.group(1) if m else None


def build_term_chapter_map(vocab_data: dict) -> dict[str, str]:
    """
    Build a lowercased-term → chapter_number map from parsed vocab data.

    Created: 2026-02-26
    """
    mapping = {}
    for ch_num, ch_terms in vocab_data['chapters'].items():
        for t in ch_terms:
            mapping[t.lower()] = ch_num
    return mapping


def lookup_vocab_chapter(
    term: str,
    term_chapter_map: dict[str, str]
) -> str | None:
    """
    Look up vocab chapter for a term using exact then normalised matching.

    Created: 2026-02-26
    """
    # Exact (case-insensitive)
    if term.lower() in term_chapter_map:
        return term_chapter_map[term.lower()]
    # Normalised (punctuation stripped)
    norm = _normalise(term)
    for vt, ch in term_chapter_map.items():
        if _normalise(vt) == norm:
            return ch
    return None


# =============================================================================
# STEP 1: DELETE NOISE
# =============================================================================

def step1_delete_noise(
    cursor: sqlite3.Cursor,
    dry_run: bool
) -> tuple[int, int]:
    """
    Delete potential_noise and high_priority_review occurrences.

    Returns (occurrences_deleted, orphan_concepts_deleted).

    Created: 2026-02-26
    """
    cursor.execute("""
        SELECT COUNT(*) FROM occurrences
        WHERE validation_status IN ('potential_noise', 'high_priority_review')
    """)
    count = cursor.fetchone()[0]

    if dry_run:
        return count, 0

    cursor.execute("""
        DELETE FROM occurrences
        WHERE validation_status IN ('potential_noise', 'high_priority_review')
    """)
    deleted = cursor.rowcount
    orphans = cleanup_orphan_concepts(cursor)
    return deleted, orphans


# =============================================================================
# STEP 2: PROMOTE CONFIRMED_WITH_FLAG
# =============================================================================

def step2_promote_flagged(
    cursor: sqlite3.Cursor,
    dry_run: bool
) -> int:
    """
    Promote confirmed_with_flag → confirmed.

    Vocab list membership overrides Stage 1 noise flags.

    Returns count updated.

    Created: 2026-02-26
    """
    cursor.execute("""
        SELECT COUNT(*) FROM occurrences
        WHERE validation_status = 'confirmed_with_flag'
    """)
    count = cursor.fetchone()[0]

    if dry_run:
        return count

    cursor.execute("""
        UPDATE occurrences SET validation_status = 'confirmed'
        WHERE validation_status = 'confirmed_with_flag'
    """)
    return cursor.rowcount


# =============================================================================
# STEP 3: UPDATE CHAPTERS FROM VOCAB LISTS
# =============================================================================

def step3_update_chapters(
    cursor: sqlite3.Cursor,
    dry_run: bool
) -> dict:
    """
    Populate NULL/empty chapter fields from vocab list chapter structure.

    Policy:
    - NULL or empty chapter → fill with vocab chapter number string ('1', '2' …)
    - Existing chapter where number matches vocab → leave unchanged
    - Existing chapter where number conflicts with vocab → log only; do NOT
      overwrite a full title like '1. The Roman Empire' with bare '2'

    Returns dict:
        filled    : occurrences where NULL chapter was populated
        conflicts : occurrences where DB chapter number ≠ vocab chapter (logged)
        no_vocab  : units skipped due to missing vocab list

    Created: 2026-02-26
    """
    units = get_all_units(cursor)
    counts = {'filled': 0, 'conflicts': 0, 'no_vocab': 0}

    for unit_meta in units:
        source_path = unit_meta['source_path']
        if not source_path:
            counts['no_vocab'] += 1
            continue

        vocab_path = find_vocab_list(source_path)
        if not vocab_path:
            counts['no_vocab'] += 1
            continue

        try:
            vocab_data = parse_vocab_docx(vocab_path)
        except Exception as e:
            print(f"  [WARN] Vocab parse error for {unit_meta['unit']}: {e}")
            counts['no_vocab'] += 1
            continue

        term_chapter_map = build_term_chapter_map(vocab_data)

        occurrences = get_unit_occurrences(
            cursor,
            unit_meta['subject'], unit_meta['year'],
            unit_meta['term'], unit_meta['unit']
        )

        for occ in occurrences:
            vocab_chapter = lookup_vocab_chapter(occ['concept_term'], term_chapter_map)
            if not vocab_chapter:
                continue  # Term not in vocab list for this unit; skip

            db_chapter_num = chapter_number_from_db_string(occ['chapter'])

            if not occ['chapter'] or not occ['chapter'].strip():
                # NULL or empty — fill in vocab chapter number
                if not dry_run:
                    cursor.execute(
                        "UPDATE occurrences SET chapter=? WHERE occurrence_id=?",
                        (str(vocab_chapter), occ['occurrence_id'])
                    )
                counts['filled'] += 1

            elif db_chapter_num and db_chapter_num != vocab_chapter:
                # Conflict — log but do not overwrite
                counts['conflicts'] += 1
                print(
                    f"  [CHAPTER CONFLICT] '{occ['concept_term']}' | "
                    f"{unit_meta['subject']} Y{unit_meta['year']} "
                    f"{unit_meta['term']} | "
                    f"DB='{occ['chapter']}' vocab_chapter='{vocab_chapter}'"
                )

    return counts


# =============================================================================
# STEP 4: RECOVER MISSED VOCAB TERMS
# =============================================================================

def step4_recover_missed(
    cursor: sqlite3.Cursor,
    dry_run: bool
) -> dict:
    """
    Find vocab list terms absent from the DB and recover them via PPTX search.

    For each unit:
    1. Load vocab list terms.
    2. Find terms not matched in current DB occurrences.
    3. Search full PPTX text for each missed term.
    4. Insert confirmed occurrence (is_introduction=0) if found.
    5. Log terms not appearing in PPTX (supplementary vocab; skip).

    Returns dict:
        recovered  : occurrences inserted
        not_found  : terms absent from PPTX text
        no_source  : units skipped due to missing PPTX or vocab list

    Created: 2026-02-26
    """
    units = get_all_units(cursor)
    counts = {'recovered': 0, 'not_found': 0, 'no_source': 0}
    total_units = len(units)

    for i, unit_meta in enumerate(units, 1):
        source_path = unit_meta['source_path']
        subject = unit_meta['subject']
        year = unit_meta['year']
        term = unit_meta['term']
        unit = unit_meta['unit']

        if i % 5 == 0 or i == total_units:
            print(f"  Step 4: unit {i}/{total_units}...")

        # Require both PPTX on disk and vocab list
        if not source_path or not Path(source_path).exists():
            counts['no_source'] += 1
            continue

        vocab_path = find_vocab_list(source_path)
        if not vocab_path:
            counts['no_source'] += 1
            continue

        try:
            vocab_data = parse_vocab_docx(vocab_path)
        except Exception as e:
            print(f"  [WARN] Vocab parse error for {unit}: {e}")
            counts['no_source'] += 1
            continue

        # Current DB terms for this unit
        occurrences = get_unit_occurrences(cursor, subject, year, term, unit)
        db_terms = [occ['concept_term'] for occ in occurrences]

        # Find vocab terms not in DB
        missed = []
        for vocab_term in vocab_data['all_terms']:
            result = match_term(vocab_term, db_terms)
            if not result['matched']:
                chapter_num = None
                for ch_num, ch_terms in vocab_data['chapters'].items():
                    if vocab_term in ch_terms:
                        chapter_num = ch_num
                        break
                missed.append({'term': vocab_term, 'chapter': chapter_num})

        if not missed:
            continue

        vocab_source = Path(vocab_path).name

        for missed_item in missed:
            term_text = missed_item['term']
            search_result = search_term_in_pptx(source_path, term_text)

            if 'error' in search_result:
                print(
                    f"  [ERROR] PPTX search failed for '{term_text}' "
                    f"in {unit}: {search_result['error']}"
                )
                counts['not_found'] += 1
                continue

            if not search_result['found']:
                counts['not_found'] += 1
                continue

            # Term found unbolded in PPTX — recover as confirmed occurrence
            slide_number = search_result['slides'][0] if search_result['slides'] else None
            term_in_context = search_result['first_context'] or None
            chapter = (
                str(missed_item['chapter']) if missed_item['chapter'] else None
            )

            counts['recovered'] += 1

            if dry_run:
                continue

            concept_id = get_or_create_concept(cursor, term_text)

            # Idempotent check — avoid duplicating an existing occurrence
            cursor.execute("""
                SELECT occurrence_id FROM occurrences
                WHERE concept_id=? AND subject=? AND year=? AND term=? AND unit=?
                AND (
                    slide_number IS ? OR
                    (slide_number IS NULL AND ? IS NULL)
                )
            """, (
                concept_id, subject, year, term, unit,
                slide_number, slide_number
            ))
            if cursor.fetchone():
                counts['recovered'] -= 1  # Already exists; don't double-count
                continue

            cursor.execute("""
                INSERT INTO occurrences (
                    concept_id, subject, year, term, unit, chapter,
                    slide_number, is_introduction, term_in_context, source_path,
                    needs_review, review_reason,
                    validation_status, vocab_confidence,
                    vocab_match_type, vocab_source
                ) VALUES (
                    ?, ?, ?, ?, ?, ?,
                    ?, 0,
                    ?, ?,
                    0, NULL,
                    'confirmed', 1.0,
                    'vocab_first_recovery', ?
                )
            """, (
                concept_id, subject, year, term, unit, chapter,
                slide_number, term_in_context, source_path,
                vocab_source
            ))

    return counts


# =============================================================================
# MAIN
# =============================================================================

def main() -> int:
    parser = argparse.ArgumentParser(
        description='Bring database to vocab-first state (4-step cleanup).'
    )
    parser.add_argument(
        '--dry-run', action='store_true',
        help='Print counts only; do not modify the database.'
    )
    parser.add_argument(
        '--skip-promote', action='store_true',
        help='Skip Step 2 (confirmed_with_flag promotion).'
    )
    args = parser.parse_args()

    project_root = Path(__file__).parent.parent
    db_path = project_root / "db" / "owl_knowledge_map.db"

    if not db_path.exists():
        print(f"ERROR: Database not found at {db_path}")
        return 1

    if args.dry_run:
        print("DRY RUN — no changes will be made to the database.")
        print()

    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # ------------------------------------------------------------------
    # Step 1
    # ------------------------------------------------------------------
    print("Step 1: Deleting noise occurrences...")
    deleted, orphans = step1_delete_noise(cursor, args.dry_run)
    if args.dry_run:
        print(f"  Would delete {deleted} occurrences "
              f"(potential_noise + high_priority_review)")
    else:
        print(f"  Deleted {deleted} occurrences")
        print(f"  Cleaned {orphans} orphan concepts")
    print()

    # ------------------------------------------------------------------
    # Step 2
    # ------------------------------------------------------------------
    promoted = 0
    if args.skip_promote:
        print("Step 2: Skipped (--skip-promote).")
    else:
        print("Step 2: Promoting confirmed_with_flag → confirmed...")
        promoted = step2_promote_flagged(cursor, args.dry_run)
        if args.dry_run:
            print(f"  Would promote {promoted} occurrences")
        else:
            print(f"  Promoted {promoted} occurrences")
    print()

    # ------------------------------------------------------------------
    # Step 3
    # ------------------------------------------------------------------
    print("Step 3: Updating chapters from vocab list structure...")
    ch_counts = step3_update_chapters(cursor, args.dry_run)
    if args.dry_run:
        print(f"  Would fill {ch_counts['filled']} NULL chapters")
    else:
        print(f"  Filled {ch_counts['filled']} NULL chapters")
    if ch_counts['conflicts'] > 0:
        print(f"  Chapter conflicts logged (not changed): {ch_counts['conflicts']}")
    if ch_counts['no_vocab'] > 0:
        print(f"  Units without vocab list: {ch_counts['no_vocab']}")
    print()

    # ------------------------------------------------------------------
    # Step 4
    # ------------------------------------------------------------------
    print("Step 4: Recovering missed vocab terms from PPTX text...")
    rec_counts = step4_recover_missed(cursor, args.dry_run)
    if args.dry_run:
        print(f"  Would recover {rec_counts['recovered']} terms (found unbolded in PPTX)")
        print(f"  Not found in PPTX: {rec_counts['not_found']} (supplementary vocab)")
    else:
        print(f"  Recovered {rec_counts['recovered']} terms")
        print(f"  Not found in PPTX: {rec_counts['not_found']} (supplementary vocab, skipped)")
    if rec_counts['no_source'] > 0:
        print(f"  Units without PPTX/vocab: {rec_counts['no_source']}")
    print()

    # ------------------------------------------------------------------
    # Commit or rollback
    # ------------------------------------------------------------------
    if args.dry_run:
        conn.rollback()
        print("DRY RUN complete — no changes made.")
    else:
        conn.commit()
        print("All changes committed.")
        print()
        print("Next steps:")
        print("  Verify:  sqlite3 db/owl_knowledge_map.db")
        print('  Query:   SELECT validation_status, COUNT(*) FROM occurrences')
        print('             GROUP BY validation_status;')
        print("  Regenerate audit (now much smaller):")
        print("           python src/audit_terms.py")
    conn.close()

    # ------------------------------------------------------------------
    # Summary
    # ------------------------------------------------------------------
    print()
    print("=" * 60)
    print("VOCAB-FIRST CLEANUP SUMMARY")
    print("=" * 60)
    label_del = "Would delete" if args.dry_run else "Deleted"
    label_prm = "Would promote" if args.dry_run else "Promoted"
    label_fil = "Would fill" if args.dry_run else "Filled"
    label_rec = "Would recover" if args.dry_run else "Recovered"

    print(f"{label_del} occurrences (noise):      {deleted}")
    if not args.dry_run:
        print(f"Orphan concepts cleaned:           {orphans}")
    if args.skip_promote:
        print(f"{label_prm} (flagged):               skipped")
    else:
        print(f"{label_prm} (flagged):               {promoted}")
    print(f"{label_fil} NULL chapters:             {ch_counts['filled']}")
    print(f"Chapter conflicts (logged only):   {ch_counts['conflicts']}")
    print(f"{label_rec} missed terms:             {rec_counts['recovered']}")
    print(f"Missed terms absent from PPTX:     {rec_counts['not_found']}")
    print("=" * 60)

    return 0


if __name__ == "__main__":
    exit(main())
