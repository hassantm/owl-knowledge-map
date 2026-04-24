#!/usr/bin/env python3
"""
Story Icon Hash Finder

Scans a PPTX file and prints the MD5 hash and dimensions of every picture
shape found. Run this on a lesson file that contains a known story slide,
identify which hash corresponds to the story icon, then pass that hash to
content_ingestion.py via --story-icon-hash.

Usage:
    python find_story_icon.py /path/to/lesson.pptx
    python find_story_icon.py /path/to/lesson.pptx --slide 3
"""

import argparse
import hashlib
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


def scan_pptx(pptx_path: str, slide_filter: int | None = None):
    prs = Presentation(pptx_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_filter is not None and slide_num != slide_filter:
            continue

        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pictures:
            continue

        print(f"\nSlide {slide_num} — {len(pictures)} picture(s):")
        for shape in pictures:
            h    = hashlib.md5(shape.image.blob).hexdigest()
            w_cm = round(Emu(shape.width).cm,  1)
            h_cm = round(Emu(shape.height).cm, 1)
            fmt  = shape.image.ext
            print(f"  {h}  {w_cm}cm × {h_cm}cm  [{fmt}]  name={shape.name!r}")


def main():
    parser = argparse.ArgumentParser(
        description="List picture shapes and MD5 hashes in a PPTX file",
    )
    parser.add_argument("pptx", help="Path to PPTX file")
    parser.add_argument("--slide", type=int, help="Only show this slide number")
    args = parser.parse_args()

    if not Path(args.pptx).exists():
        print(f"ERROR: file not found: {args.pptx}")
        return 1

    scan_pptx(args.pptx, args.slide)
    return 0


if __name__ == "__main__":
    sys.exit(main())
