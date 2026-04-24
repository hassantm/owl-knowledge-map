#!/usr/bin/env python3
"""
Content Ingestion for Story Pack Generator

Extracts and tokenises curriculum content from PPTX files, storing
structured JSONB in the units table. Two document types are handled:

  lesson  — teacher-facing PPTX with speaker notes, story slide markers,
             source page references, and vocabulary cues
  booklet — pupil-facing PPTX used as the factual source for story generation

Token counts are computed at ingestion time using the Anthropic token-counting
API so context assembly never needs to call it again at generation time.

Usage (standalone):
    python content_ingestion.py --unit-id 35 --lesson /path/to/lesson.pptx
    python content_ingestion.py --unit-id 35 --booklet /path/to/booklet.pptx
    python content_ingestion.py --unit-id 35 --lesson /path/to/lesson.pptx \
                                              --booklet /path/to/booklet.pptx
    python content_ingestion.py --unit-id 35 --lesson /path/to/lesson.pptx --dry-run
"""

import argparse
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

import anthropic
import psycopg2.extras
from dotenv import load_dotenv
from pptx import Presentation

from db import get_connection

load_dotenv()

client = anthropic.Anthropic()

# Phrases in slide title or speaker notes that mark a story slide
STORY_SLIDE_PATTERNS = [
    r"listen(?:ing)? to a story",
    r"tell(?:ing)? the story",
    r"tell(?:ing)? a story",
    r"hear(?:ing)? a story",
    r"story time",
    r"let'?s listen",
    r"let'?s hear",
]

_STORY_RE = re.compile("|".join(STORY_SLIDE_PATTERNS), re.IGNORECASE)

# Matches "Page 3", "pages 5 and 6", "Pages 5, 6 and 7"
_SOURCE_PAGES_RE = re.compile(
    r"[Pp]ages?\s+((?:\d+(?:[,\s]+and\s+|\s+and\s+|\s*,\s*|\s+))*\d+)",
)

# Matches: "the word 'precinct'" or "the words 'x' and 'y'"
_VOCAB_RE = re.compile(r"the words?\s+'([^']+)'", re.IGNORECASE)

# Animation / visual instructions that should be preserved separately
_ANIMATION_RE = re.compile(
    r"(?:Remove|Use|Show|Hide|Click|Reveal|Animate)[^.]{0,300}\.",
    re.IGNORECASE,
)


def count_tokens(text: str, model: str) -> int:
    """Count tokens for a text chunk using the Anthropic token-counting API."""
    if not text.strip():
        return 0
    response = client.messages.count_tokens(
        model=model,
        messages=[{"role": "user", "content": text}],
    )
    return response.input_tokens


def _normalise_text(text: str) -> str:
    """Normalise curly/smart quotes to plain ASCII so regex matches reliably."""
    return (text
            .replace("\u2018", "'").replace("\u2019", "'")
            .replace("\u201c", '"').replace("\u201d", '"'))


def detect_story_slide(title: str, notes: str, story_icon_hashes: set[str] | None = None,
                       slide=None) -> bool:
    combined = _normalise_text((title or "") + " " + (notes or ""))
    if _STORY_RE.search(combined):
        return True
    if story_icon_hashes and slide is not None:
        return _slide_has_story_icon(slide, story_icon_hashes)
    return False


def _slide_has_story_icon(slide, story_icon_hashes: set[str]) -> bool:
    """Return True if any picture shape on the slide matches a known story icon hash."""
    import hashlib
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            h = hashlib.md5(shape.image.blob).hexdigest()
            if h in story_icon_hashes:
                return True
    return False


def parse_source_pages(notes: str) -> list[int]:
    pages = []
    for m in _SOURCE_PAGES_RE.finditer(notes or ""):
        pages.extend(int(x) for x in re.findall(r"\d+", m.group(1)))
    return sorted(set(pages))


def parse_vocabulary_notes(notes: str) -> list[str]:
    return _VOCAB_RE.findall(notes or "")


def parse_animation_notes(notes: str) -> str | None:
    m = _ANIMATION_RE.search(notes or "")
    return m.group(0).strip() if m else None


def _slide_title(slide) -> str:
    """Return the title text of a slide, or empty string if none."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def _slide_body_text(slide) -> str:
    """Return all non-title text from a slide, joined by newlines."""
    parts = []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _slide_notes(slide) -> str:
    """Return the speaker notes text, or empty string."""
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def _extract_slides_from_pptx(pptx_path: str, model: str, start_index: int = 1,
                               story_icon_hashes: set[str] | None = None) -> tuple[dict, int]:
    """
    Extract slides from a single PPTX file, numbering from start_index.
    Returns (slides dict, total_tokens).
    """
    prs = Presentation(pptx_path)
    slides = {}
    total_tokens = 0

    for i, slide in enumerate(prs.slides, start=start_index):
        title = _slide_title(slide)
        body  = _slide_body_text(slide)
        notes = _slide_notes(slide)

        text = "\n".join(filter(None, [title, body]))
        combined = text + " " + notes

        token_count = count_tokens(combined, model)
        total_tokens += token_count

        is_story = detect_story_slide(title, notes, story_icon_hashes, slide)

        entry: dict = {
            "text":        text,
            "notes":       notes,
            "token_count": token_count,
            "story_slide": is_story,
        }

        if is_story:
            entry["source_pages"]     = parse_source_pages(notes)
            entry["vocabulary_notes"] = parse_vocabulary_notes(notes)

        anim = parse_animation_notes(notes)
        if anim:
            entry["animation_notes"] = anim

        slides[str(i)] = entry

    return slides, total_tokens


def extract_lesson_content(pptx_path: str, model: str,
                           story_icon_hashes: set[str] | None = None) -> dict:
    """
    Extract lesson content into structured JSONB dict.

    pptx_path may be a single .pptx file or a directory containing multiple
    .pptx files. If a directory, all .pptx files are processed in alphabetical
    order with continuous slide numbering across files.

    Each slide entry includes:
      text             — visible slide text (title + body, newline-separated)
      notes            — speaker notes
      token_count      — token count for text + notes combined
      story_slide      — True if this is a story-telling slide
      source_pages     — booklet page numbers referenced in notes (story slides only)
      vocabulary_notes — terms flagged for embedding naturally (story slides only)
      animation_notes  — visual/animation instruction from notes (if present)
    """
    path = Path(pptx_path)
    if path.is_dir():
        pptx_files = sorted(path.glob("*.pptx"))
        if not pptx_files:
            raise ValueError(f"No .pptx files found in {pptx_path}")
    else:
        pptx_files = [path]

    slides = {}
    total_tokens = 0
    next_index = 1

    for pptx_file in pptx_files:
        file_slides, file_tokens = _extract_slides_from_pptx(
            str(pptx_file), model, next_index, story_icon_hashes
        )
        slides.update(file_slides)
        total_tokens += file_tokens
        next_index += len(file_slides)

    story_count = sum(1 for s in slides.values() if s["story_slide"])

    return {
        "document_type":     "lesson",
        "total_token_count": total_tokens,
        "tokenizer_model":   model,
        "extracted_at":      datetime.now(timezone.utc).isoformat(),
        "slide_count":       len(slides),
        "story_slide_count": story_count,
        "slides":            slides,
    }


def extract_booklet_content(pptx_path: str, model: str) -> dict:
    """
    Extract booklet PPTX into structured JSONB dict.

    Booklet slides are treated as numbered pages. Only visible text is
    captured — no speaker note analysis needed for the booklet.
    """
    prs = Presentation(pptx_path)
    pages = {}
    total_tokens = 0

    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body  = _slide_body_text(slide)
        text  = "\n".join(filter(None, [title, body]))

        token_count = count_tokens(text, model)
        total_tokens += token_count

        pages[str(i)] = {
            "text":        text,
            "token_count": token_count,
        }

    return {
        "document_type":     "booklet",
        "total_token_count": total_tokens,
        "tokenizer_model":   model,
        "extracted_at":      datetime.now(timezone.utc).isoformat(),
        "page_count":        len(pages),
        "pages":             pages,
    }


def _write_unit_content(unit_id: int, content: dict, document_type: str, model: str):
    col = f"{document_type}_content"
    conn = get_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(f"""
                UPDATE units SET
                    {col}               = %s,
                    {col}_stale         = false,
                    {col}_extracted_at  = now(),
                    {col}_model         = %s
                WHERE unit_id = %s
            """, (psycopg2.extras.Json(content), model, unit_id))
        conn.commit()
    finally:
        conn.close()


def ingest_lesson(pptx_path: str, unit_id: int, model: str = "claude-sonnet-4-6",
                  dry_run: bool = False,
                  story_icon_hashes: set[str] | None = None) -> dict:
    """
    Extract and store lesson content for a unit.
    Returns the content dict (whether or not dry_run).
    """
    print(f"Extracting lesson: {Path(pptx_path).name}")
    content = extract_lesson_content(pptx_path, model, story_icon_hashes)

    story_slides = content["story_slide_count"]
    print(f"  {content['slide_count']} slides, {story_slides} story slide(s), "
          f"{content['total_token_count']} tokens")

    if dry_run:
        print("  [dry-run] skipping database write")
        return content

    _write_unit_content(unit_id, content, "lesson", model)
    print(f"  Written to units.lesson_content (unit_id={unit_id})")
    return content


def ingest_booklet(pptx_path: str, unit_id: int, model: str = "claude-sonnet-4-6",
                   dry_run: bool = False) -> dict:
    """
    Extract and store booklet content for a unit.
    Returns the content dict (whether or not dry_run).
    """
    print(f"Extracting booklet: {Path(pptx_path).name}")
    content = extract_booklet_content(pptx_path, model)

    print(f"  {content['page_count']} pages, {content['total_token_count']} tokens")

    if dry_run:
        print("  [dry-run] skipping database write")
        return content

    _write_unit_content(unit_id, content, "booklet", model)
    print(f"  Written to units.booklet_content (unit_id={unit_id})")
    return content


def mark_stale(unit_id: int, document_type: str):
    """Mark a unit's content as stale (e.g. after the source file changes)."""
    if document_type not in ("lesson", "booklet"):
        raise ValueError(f"document_type must be 'lesson' or 'booklet', got {document_type!r}")
    col = f"{document_type}_content_stale"
    conn = get_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(f"UPDATE units SET {col} = true WHERE unit_id = %s", (unit_id,))
        conn.commit()
    finally:
        conn.close()


def get_content_status(unit_id: int) -> dict:
    """Return content metadata for a unit (extracted_at, model, stale flags)."""
    conn = get_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT
                    lesson_content_stale,
                    lesson_content_extracted_at,
                    lesson_content_model,
                    lesson_content IS NOT NULL      AS has_lesson,
                    booklet_content_stale,
                    booklet_content_extracted_at,
                    booklet_content_model,
                    booklet_content IS NOT NULL     AS has_booklet
                FROM units WHERE unit_id = %s
            """, (unit_id,))
            row = cur.fetchone()
    finally:
        conn.close()
    return dict(row) if row else {}


def main():
    parser = argparse.ArgumentParser(
        description="Ingest lesson and/or booklet PPTX content for a curriculum unit",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Ingest lesson from a single PPTX
  python content_ingestion.py --unit-id 35 --lesson /path/to/Y5_Baghdad_Lesson.pptx

  # Ingest lesson from a directory of PPTXs (processed alphabetically)
  python content_ingestion.py --unit-id 35 --lesson /path/to/Baghdad/Lessons/

  # Ingest both lesson and booklet in one pass
  python content_ingestion.py --unit-id 35 \\
    --lesson  /path/to/Baghdad/Lessons/ \\
    --booklet /path/to/Y5_Baghdad_Booklet.pptx

  # Dry run — extract and print stats without writing to DB
  python content_ingestion.py --unit-id 35 --lesson /path/to/lesson.pptx --dry-run

  # Check what's already ingested for a unit
  python content_ingestion.py --unit-id 35 --status
        """,
    )
    parser.add_argument("--unit-id",  type=int, required=True, help="units.unit_id")
    parser.add_argument("--lesson",   help="Path to lesson PPTX, or directory containing multiple lesson PPTXs")
    parser.add_argument("--booklet",  help="Path to booklet PPTX")
    parser.add_argument("--model",    default="claude-sonnet-4-6",
                        help="Tokenizer model (default: claude-sonnet-4-6)")
    parser.add_argument("--dry-run",  action="store_true",
                        help="Extract and print stats without writing to database")
    parser.add_argument("--status",   action="store_true",
                        help="Print current content status for the unit and exit")
    parser.add_argument("--story-icon-hash", action="append", dest="story_icon_hashes",
                        metavar="MD5",
                        help="MD5 hash of the story icon image (can be repeated). "
                             "Use find_story_icon.py to extract hashes from a known story slide.")
    args = parser.parse_args()

    if args.status:
        status = get_content_status(args.unit_id)
        if not status:
            print(f"No unit found with unit_id={args.unit_id}")
            return 1
        print(f"Unit {args.unit_id} content status:")
        print(f"  Lesson:  {'yes' if status['has_lesson'] else 'no'}", end="")
        if status["has_lesson"]:
            print(f"  (model={status['lesson_content_model']}, "
                  f"extracted={status['lesson_content_extracted_at']}, "
                  f"stale={status['lesson_content_stale']})", end="")
        print()
        print(f"  Booklet: {'yes' if status['has_booklet'] else 'no'}", end="")
        if status["has_booklet"]:
            print(f"  (model={status['booklet_content_model']}, "
                  f"extracted={status['booklet_content_extracted_at']}, "
                  f"stale={status['booklet_content_stale']})", end="")
        print()
        return 0

    if not args.lesson and not args.booklet:
        parser.error("Provide at least one of --lesson or --booklet")

    icon_hashes = set(args.story_icon_hashes) if args.story_icon_hashes else None

    if args.lesson:
        p = Path(args.lesson)
        if not p.exists():
            print(f"ERROR: lesson path not found: {args.lesson}")
            return 1
        if p.is_dir() and not list(p.glob("*.pptx")):
            print(f"ERROR: no .pptx files found in {args.lesson}")
            return 1
        ingest_lesson(args.lesson, args.unit_id, args.model, dry_run=args.dry_run,
                      story_icon_hashes=icon_hashes)

    if args.booklet:
        ingest_booklet(args.booklet, args.unit_id, args.model, dry_run=args.dry_run)

    return 0


if __name__ == "__main__":
    sys.exit(main())
