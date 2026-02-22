#!/usr/bin/env python3
"""
Stage 1: PPTX Bold Term Extraction (Proof of Concept)

Extracts bold text from curriculum PowerPoint booklets, treating bold as pedagogical
markers for formal concept introductions. Applies noise filtering, tracks chapters,
and outputs results for human validation.

No database writes in Stage 1—focuses on validating extraction logic.

Created: 2026-02-22
"""

import re
from pathlib import Path
from pptx import Presentation


# =============================================================================
# NOISE TERMS DICTIONARY
# =============================================================================

# 2026-02-22: Noise terms dictionary - slide labels and structural text
NOISE_TERMS = {
    # Slide structural labels
    'reason 1', 'reason 2', 'reason 3', 'reason 4', 'reason 5',
    'source 1', 'source 2', 'source 3', 'source 4', 'source 5',
    'example 1', 'example 2', 'example 3',
    'task 1', 'task 2', 'task 3',

    # Generic headings (only when standalone)
    'the', 'a', 'an',

    # Add more patterns as discovered during corpus processing
}


# =============================================================================
# LAYER 1: FILTERING & CLEANING
# =============================================================================

def is_noise(text: str) -> bool:
    """
    Filter out non-conceptual bold text patterns.

    Filters:
    - "Page N" patterns (table of contents artifacts)
    - Pure numeric runs like "456" or "17." (line numbers from reading scaffolds)
    - URLs (http://, https://, www.)
    - Common citation patterns (Group, Inc, Ltd, LLC, Corp, Organization, Foundation)
    - Noise terms from dictionary (Reason 1, Source 1, etc.)

    Args:
        text: Raw text from a bold run

    Returns:
        True if text should be filtered out, False if it's a valid term
    """
    text = text.strip()

    # Filter "Page N" patterns
    if re.match(r'^Page\s+\d+$', text, re.IGNORECASE):
        return True

    # Filter pure numeric runs (with optional trailing period)
    if re.match(r'^\d+\.?$', text):
        return True

    # 2026-02-22: Filter URLs
    if re.match(r'https?://', text, re.IGNORECASE):
        return True
    if re.match(r'www\.', text, re.IGNORECASE):
        return True

    # 2026-02-22: Filter common citation patterns
    if re.match(r'.+(Group|Inc|Ltd|LLC|Corp|Organization|Foundation)\.?$', text):
        return True

    # 2026-02-22: Filter noise terms from dictionary (clean first to handle trailing punctuation)
    cleaned = text.rstrip('.,;:!?')
    if cleaned.lower() in NOISE_TERMS:
        return True

    return False


def clean_term(text: str) -> str:
    """
    Clean a term by stripping trailing punctuation.

    Preserves internal punctuation in multi-word terms (e.g., "St. Paul").
    Only strips trailing: . , ; : ! ?

    Args:
        text: Raw term text

    Returns:
        Cleaned term
    """
    return text.strip().rstrip('.,;:!?')


def flag_for_review(text: str, context: str) -> tuple[bool, str]:
    """
    Determine if term needs human review and why.

    2026-02-22: Enhanced flagging with specific reasons

    Args:
        text: Cleaned term text
        context: Full paragraph text surrounding the term

    Returns:
        (needs_review: bool, reason: str)
    """
    reasons = []

    # Short terms
    if len(text) < 5:
        reasons.append('short_term')

    # Very short context suggests heading-only
    if len(context) < 20:
        reasons.append('potential_heading')

    # Single word in all caps (except valid proper nouns)
    if text.isupper() and ' ' not in text and len(text) > 1:
        reasons.append('all_caps')

    # Ends with colon (likely a heading)
    if text.endswith(':'):
        reasons.append('heading_marker')

    if reasons:
        return True, ', '.join(reasons)
    return False, None


def detect_chapter(text: str) -> str:
    """
    Detect chapter headings from paragraph text.

    Pattern: "1. Chapter title" or "2. Another chapter"

    Args:
        text: Full paragraph text

    Returns:
        Chapter heading if detected, None otherwise
    """
    text = text.strip()
    match = re.match(r'^(\d+\.\s+.+)$', text)
    if match:
        return match.group(1)
    return None


# =============================================================================
# LAYER 2: EXTRACTION
# =============================================================================

def extract_bold_runs(pptx_path: str) -> dict:
    """
    Extract bold terms from a PowerPoint file with chapter tracking.

    Core extraction algorithm:
    1. Iterate through slides, shapes, paragraphs, runs
    2. Detect chapter headings from paragraph text
    3. Extract runs where font.bold == True (explicit bold only)
    4. Apply noise filtering and cleaning
    5. Track slide number, chapter, and paragraph context

    Args:
        pptx_path: Path to PowerPoint file

    Returns:
        dict with keys:
            - terms: List of extracted term dicts
            - total_slides: Number of slides processed
            - errors: List of error messages (if any)
    """
    results = {
        'terms': [],
        'total_slides': 0,
        'errors': []
    }

    current_chapter = None
    in_credits_section = False  # 2026-02-22: Track Picture Credits section

    try:
        prs = Presentation(pptx_path)
        results['total_slides'] = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                # Skip shapes without text frames
                if not hasattr(shape, 'text_frame'):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    # Get full paragraph text for context and chapter detection
                    para_text = paragraph.text.strip()

                    # 2026-02-22: Check if we've entered Picture Credits section
                    if 'picture credit' in para_text.lower():
                        in_credits_section = True

                    # Skip all bold terms if in credits section
                    if in_credits_section:
                        continue

                    # Check for chapter heading (scans ALL text, not just bold)
                    chapter_heading = detect_chapter(para_text)
                    if chapter_heading:
                        current_chapter = chapter_heading

                    # Extract bold runs
                    for run in paragraph.runs:
                        # Only capture explicit bold (True), not inherited (None)
                        if run.font.bold is not True:
                            continue

                        run_text = run.text.strip()

                        # Skip empty runs
                        if not run_text:
                            continue

                        # Apply noise filter
                        if is_noise(run_text):
                            continue

                        # 2026-02-22: Clean and flag with enhanced review logic
                        cleaned = clean_term(run_text)
                        needs_review, review_reason = flag_for_review(cleaned, para_text)

                        # Store extracted term with metadata
                        results['terms'].append({
                            'term': cleaned,
                            'slide': slide_num,
                            'chapter': current_chapter,
                            'context': para_text,
                            'flagged': needs_review,
                            'review_reason': review_reason
                        })

    except Exception as e:
        results['errors'].append(f"Extraction error: {str(e)}")

    return results


# =============================================================================
# LAYER 3: OUTPUT
# =============================================================================

def format_output(results: dict, filename: str) -> None:
    """
    Print structured console report of extraction results.

    Three-section format:
    1. Summary statistics
    2. Terms grouped by slide and chapter
    3. Flagged terms for review

    Args:
        results: Extraction results from extract_bold_runs()
        filename: Source filename for display
    """
    terms = results['terms']
    flagged_terms = [t for t in terms if t['flagged']]

    # Section 1: Summary
    print("=" * 60)
    print("=== EXTRACTION REPORT ===")
    print("=" * 60)
    print(f"File: {filename}")
    print(f"Slides processed: {results['total_slides']}")
    print(f"Terms extracted: {len(terms)}")
    print(f"Flagged for review: {len(flagged_terms)}")

    if results['errors']:
        print("\nERRORS:")
        for error in results['errors']:
            print(f"  - {error}")

    # Section 2: Extracted Terms
    print("\n" + "=" * 60)
    print("=== EXTRACTED TERMS ===")
    print("=" * 60)

    # Group by slide for readability
    current_slide = None
    for term_data in terms:
        slide_num = term_data['slide']
        chapter = term_data['chapter']
        term = term_data['term']
        context = term_data['context']
        flagged = term_data['flagged']

        # Print slide/chapter header when it changes
        if slide_num != current_slide:
            current_slide = slide_num
            chapter_display = f" | {chapter}" if chapter else ""
            print(f"\nSlide {slide_num}{chapter_display}")

        # 2026-02-22: Print term with flag and reason if applicable
        flag_marker = f"[FLAGGED: {term_data.get('review_reason', 'SHORT')}] " if flagged else ""
        print(f'  {flag_marker}"{term}"')

        # Print context (truncate if very long)
        context_display = context if len(context) <= 100 else context[:97] + "..."
        print(f'    Context: {context_display}')

    # Section 3: Flagged Terms Summary
    if flagged_terms:
        print("\n" + "=" * 60)
        print("=== FLAGGED TERMS (for review) ===")
        print("=" * 60)
        for i, term_data in enumerate(flagged_terms, start=1):
            term = term_data['term']
            slide = term_data['slide']
            reason = term_data.get('review_reason', 'unknown')
            print(f'{i}. "{term}" ({len(term)} chars) - Slide {slide} - Reason: {reason}')


def main():
    """
    Main execution: Extract bold terms from sample PPTX and print report.
    """
    # Path to sample file
    project_root = Path(__file__).parent.parent
    sample_file = project_root / "data" / "sample" / "Y4 Spring 2 Christianity in 3 empires Booklet.pptx"

    # Verify file exists
    if not sample_file.exists():
        print(f"ERROR: Sample file not found at {sample_file}")
        return

    print(f"Processing: {sample_file.name}")
    print("Extracting bold terms...\n")

    # Extract
    results = extract_bold_runs(str(sample_file))

    # Output
    format_output(results, sample_file.name)


if __name__ == "__main__":
    main()
