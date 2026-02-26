#!/usr/bin/env python3
"""
Audit Enrichment Script

Reads output/term_audit.csv and enriches it with:
- occurrence_id for noise/high_priority rows (needed for delete operations)
- PPTX text search results for missed_from_extraction rows:
    appears_unbolded, unbolded_slides, unbolded_context
- blank decision column for reviewer to fill in

Writes output/term_audit_enriched.csv.

Valid decisions (document for reviewer):
  missed_from_extraction : add (only if appears_unbolded=True), skip
  potential_noise        : keep, delete, skip
  high_priority_review   : keep, delete, skip
  (blank = not yet reviewed, treated as skip by apply_audit_decisions.py)

Created: 2026-02-24
"""

import csv
import re
import sqlite3
from pathlib import Path

from pptx import Presentation


# =============================================================================
# DATABASE QUERIES
# =============================================================================

def get_source_path(db_path: Path, subject: str, year: str, term: str, unit: str) -> str | None:
    """
    Get source PPTX path for a unit from the database.

    Created: 2026-02-24
    """
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute(
        """SELECT DISTINCT source_path FROM occurrences
           WHERE subject=? AND year=? AND term=? AND unit=? LIMIT 1""",
        (subject, int(year), term, unit)
    )
    row = cursor.fetchone()
    conn.close()
    return row[0] if row else None


def get_occurrence_id(
    db_path: Path,
    subject: str,
    year: str,
    term: str,
    unit: str,
    concept_term: str,
    slide_number
) -> int | None:
    """
    Look up occurrence_id for a specific term occurrence.

    Matches on subject, year, term, unit, concept term and slide number.
    Falls back to first match without slide number if slide is missing.

    Created: 2026-02-24
    """
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    try:
        slide_int = int(slide_number) if slide_number else None
    except (ValueError, TypeError):
        slide_int = None

    if slide_int is not None:
        cursor.execute(
            """
            SELECT o.occurrence_id FROM occurrences o
            JOIN concepts c ON o.concept_id = c.concept_id
            WHERE o.subject=? AND o.year=? AND o.term=? AND o.unit=?
            AND c.term=? AND o.slide_number=?
            LIMIT 1
            """,
            (subject, int(year), term, unit, concept_term, slide_int)
        )
    else:
        cursor.execute(
            """
            SELECT o.occurrence_id FROM occurrences o
            JOIN concepts c ON o.concept_id = c.concept_id
            WHERE o.subject=? AND o.year=? AND o.term=? AND o.unit=?
            AND c.term=?
            LIMIT 1
            """,
            (subject, int(year), term, unit, concept_term)
        )

    row = cursor.fetchone()
    conn.close()
    return row[0] if row else None


# =============================================================================
# PPTX TEXT SEARCH
# =============================================================================

def search_term_in_pptx(pptx_path: str, term: str) -> dict:
    """
    Search all text (including unbolded) in a PPTX for a term.

    Uses word-boundary regex, case-insensitive.
    Iterates slides > shapes > text frames > paragraphs > assembled paragraph text.
    This catches terms split across run boundaries.

    Returns:
        dict with:
        - found: bool
        - slides: sorted list of int slide numbers (1-based)
        - first_context: str — paragraph text of first match
        - error: str (only present on failure)

    Created: 2026-02-24
    """
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {'found': False, 'slides': [], 'first_context': '', 'error': str(e)}

    pattern = re.compile(
        r'(?<![a-zA-Z])' + re.escape(term) + r'(?![a-zA-Z])',
        re.IGNORECASE
    )

    matching_slides = []
    first_context = ''

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_matched = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # Assemble full paragraph text across all runs
                para_text = ''.join(run.text for run in para.runs)
                if para_text and pattern.search(para_text):
                    slide_matched = True
                    if not first_context:
                        first_context = para_text.strip()
        if slide_matched:
            matching_slides.append(slide_idx)

    return {
        'found': bool(matching_slides),
        'slides': sorted(matching_slides),
        'first_context': first_context
    }


# =============================================================================
# MAIN ENRICHMENT
# =============================================================================

def enrich_audit(input_path: Path, output_path: Path, db_path: Path) -> None:
    """
    Read term_audit.csv, add enrichment columns, write term_audit_enriched.csv.

    New columns:
        occurrence_id    — for noise/high_priority rows
        appears_unbolded — True / False / No source found (missed rows only)
        unbolded_slides  — comma-separated slide numbers (missed rows only)
        unbolded_context — first matching paragraph (missed rows only)
        decision         — blank; reviewer fills in

    Created: 2026-02-24
    """
    with open(input_path, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        rows = list(reader)

    print(f"Loaded {len(rows)} rows from {input_path.name}")

    # Cache source paths to avoid repeated DB queries for same unit
    source_path_cache: dict[tuple, str | None] = {}

    counts = {
        'missed_found': 0,
        'missed_not_found': 0,
        'missed_no_source': 0,
        'noise_id_found': 0,
        'noise_id_missing': 0,
    }

    enriched_rows = []
    total = len(rows)

    for i, row in enumerate(rows):
        issue_type = row['issue_type']
        subject = row['subject']
        year = row['year']
        term = row['term_period']
        unit = row['unit']
        concept_term = row['term']
        slide = row.get('slide', '')

        # Initialise new columns with empty defaults
        row['occurrence_id'] = ''
        row['appears_unbolded'] = ''
        row['unbolded_slides'] = ''
        row['unbolded_context'] = ''
        row['decision'] = ''

        if issue_type == 'missed_from_extraction':
            cache_key = (subject, year, term, unit)
            if cache_key not in source_path_cache:
                source_path_cache[cache_key] = get_source_path(
                    db_path, subject, year, term, unit
                )

            source_path = source_path_cache[cache_key]

            if not source_path or not Path(source_path).exists():
                row['appears_unbolded'] = 'No source found'
                counts['missed_no_source'] += 1
                if not source_path:
                    print(f"  [WARN] No source_path in DB for {subject} Y{year} {term} | {unit}")
                else:
                    print(f"  [WARN] PPTX not on disk: {source_path}")
            else:
                if i % 20 == 0:
                    print(f"  Searching missed terms... {i + 1}/{total}")

                result = search_term_in_pptx(source_path, concept_term)

                if 'error' in result:
                    row['appears_unbolded'] = 'Error'
                    row['unbolded_context'] = result['error']
                    print(f"  [ERROR] PPTX read failed for '{concept_term}': {result['error']}")
                elif result['found']:
                    row['appears_unbolded'] = 'True'
                    row['unbolded_slides'] = ', '.join(str(s) for s in result['slides'])
                    row['unbolded_context'] = result['first_context']
                    counts['missed_found'] += 1
                else:
                    row['appears_unbolded'] = 'False'
                    counts['missed_not_found'] += 1

        elif issue_type in ('potential_noise', 'high_priority_review'):
            occ_id = get_occurrence_id(
                db_path, subject, year, term, unit, concept_term, slide
            )
            if occ_id is not None:
                row['occurrence_id'] = str(occ_id)
                counts['noise_id_found'] += 1
            else:
                counts['noise_id_missing'] += 1
                print(
                    f"  [WARN] occurrence_id not found: '{concept_term}' | "
                    f"{subject} Y{year} {term} | {unit} | slide={slide}"
                )

        enriched_rows.append(row)

    # Write enriched CSV
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fieldnames = [
        'issue_type', 'subject', 'year', 'term_period', 'unit',
        'chapter', 'term', 'slide', 'context', 'review_reason',
        'vocab_source', 'notes',
        'occurrence_id', 'appears_unbolded', 'unbolded_slides', 'unbolded_context',
        'decision'
    ]

    with open(output_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(enriched_rows)

    print()
    print("=" * 60)
    print("ENRICHMENT SUMMARY")
    print("=" * 60)
    print(f"Missed terms found in PPTX (unbolded): {counts['missed_found']}")
    print(f"Missed terms NOT found in PPTX:        {counts['missed_not_found']}")
    print(f"Missed terms — no source PPTX:         {counts['missed_no_source']}")
    print(f"Noise/HP rows — occurrence_id found:   {counts['noise_id_found']}")
    print(f"Noise/HP rows — occurrence_id missing: {counts['noise_id_missing']}")
    print()
    print(f"Enriched CSV written to: {output_path}")
    print()
    print("Valid decisions to fill into 'decision' column:")
    print("  missed_from_extraction : add (only if appears_unbolded=True), skip")
    print("  potential_noise        : keep, delete, skip")
    print("  high_priority_review   : keep, delete, skip")
    print("  (leave blank = skip — no action taken)")
    print("=" * 60)


# =============================================================================
# MAIN
# =============================================================================

def main() -> int:
    project_root = Path(__file__).parent.parent
    db_path = project_root / "db" / "owl_knowledge_map.db"
    input_path = project_root / "output" / "term_audit.csv"
    output_path = project_root / "output" / "term_audit_enriched.csv"

    if not db_path.exists():
        print(f"ERROR: Database not found at {db_path}")
        return 1

    if not input_path.exists():
        print(f"ERROR: Input audit CSV not found at {input_path}")
        print("Run: python src/audit_terms.py")
        return 1

    print("Starting audit enrichment...")
    print()
    enrich_audit(input_path, output_path, db_path)
    return 0


if __name__ == "__main__":
    exit(main())
